import os
import json
from pptx import Presentation
import pdfplumber

def extract_text_from_ppt(ppt_path):
    """提取PPT中的文字"""
    ppt = Presentation(ppt_path)
    text = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text.append(paragraph.text.strip())
    return "\n".join(text)

def extract_text_from_pdf(pdf_path):
    """提取PDF中的文字"""
    text = []
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text.append(page.extract_text())
    return "\n".join(text)

def save_extracted_text(output_path, file_name, content):
    """保存提取的内容为txt文件"""
    os.makedirs(output_path, exist_ok=True)
    file_path = os.path.join(output_path, f"{file_name}.txt")
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(content)

def process_files(input_folder, output_folder):
    """处理输入文件夹中的PPT和PDF文件，并保存提取后的文本"""
    extracted_files = []
    for root, _, files in os.walk(input_folder):
        for file in files:
            file_path = os.path.join(root, file)
            file_name, file_ext = os.path.splitext(file)
            
            if file_ext.lower() == ".pptx":
                print(f"Processing PPT: {file}")
                text = extract_text_from_ppt(file_path)
            elif file_ext.lower() == ".pdf":
                print(f"Processing PDF: {file}")
                text = extract_text_from_pdf(file_path)
            else:
                print(f"Skipping unsupported file for processing pdf2txt: {file}")
                continue
            
            save_extracted_text(output_folder, file_name, text)
            extracted_files.append(file_name)
    
    return sorted(extracted_files, key=str)

def merge_text_files(output_folder, merged_file_path, sorted_files):
    """合并所有提取的文本文件为一个总文件"""
    with open(merged_file_path, "w", encoding="utf-8") as merged_file:
        for file_name in sorted_files:
            file_path = os.path.join(output_folder, f"{file_name}.txt")
            if os.path.exists(file_path):
                with open(file_path, "r", encoding="utf-8") as f:
                    merged_file.write(f"### {file_name} ###\n")
                    merged_file.write(f.read())
                    merged_file.write("\n\n")
            else:
                print(f"Warning: File {file_path} not found, skipping.")

def main():
    # 设置输入和输出路径
    input_folder = "output/自然辩证法/refs"
    output_folder = "output/自然辩证法/refs/text"
    merged_file_path = "output/自然辩证法/refs/text/merged_text.txt"
    
    # 提取文件内容并保存
    sorted_files = process_files(input_folder, output_folder)
    print("文字提取完成，所有文本已保存到指定文件夹。")
    
    # 合并所有文件内容
    merge_text_files(output_folder, merged_file_path, sorted_files)
    print(f"所有文本文件已合并，保存为: {merged_file_path}")

if __name__ == "__main__":
    main()